"""
REXY FILE READER PLUGIN
Reads files from the rexy_inbox/ folder inside your Rexy project.

How to use:
1. Drop any file into rexy_inbox/ folder
2. Say "read notes.txt" or "open my resume" or "what's in report.pdf"

Supported formats:
- .txt  → plain text
- .md   → markdown (shown as plain text)
- .csv  → shows first 20 rows
- .json → pretty printed
- .pdf  → extracts text (requires pypdf: pip install pypdf)
- .pptx → extracts slide text (requires python-pptx: pip install python-pptx)
- .docx → extracts paragraphs (requires python-docx: pip install python-docx)
"""

import os
import re
import logging
from typing import Any, Dict, List, Optional

from modules.plugin_base import RexyPlugin

logger = logging.getLogger("rexy.filereader")

# Inbox folder — create this in your Rexy project root
INBOX_DIR = "rexy_inbox"

# Max characters to show in reply (avoids giant walls of text)
MAX_CHARS = 1500


class FileReaderPlugin(RexyPlugin):
    """Read files dropped into the rexy_inbox/ folder."""

    @property
    def intent_name(self) -> str:
        return "FILE_READ"

    @property
    def description(self) -> str:
        return "Read files from the rexy_inbox folder"

    @property
    def risk_level(self) -> str:
        return "low"

    @property
    def intent_examples(self) -> List[str]:
        return [
            "read notes.txt",
            "open my resume",
            "what's in report.pdf",
            "read the file called homework",
            "show me inbox files",
        ]

    # ── Main execute ──
    def execute(self, message: str, emotion: str, state: Dict[str, Any]) -> Dict[str, Any]:
        """Figure out what file the user wants and read it."""

        # Ensure inbox folder exists
        self._ensure_inbox()

        message_lower = message.lower().strip()

        # ── LIST FILES IN INBOX ──
        if re.search(r'\b(list|show|what files|what\'s in|inbox)\b', message_lower):
            return self._list_inbox()

        # ── READ A SPECIFIC FILE ──
        filename = self._extract_filename(message)

        if not filename:
            # No filename found — list inbox to help user
            return self._list_inbox(hint=True)

        return self._read_file(filename)

    # ─────────────────────────────────────────────
    # LIST INBOX
    # ─────────────────────────────────────────────
    def _list_inbox(self, hint: bool = False) -> Dict[str, Any]:
        """Show all files currently in rexy_inbox/."""
        try:
            files = os.listdir(INBOX_DIR)
            # Filter out hidden files and folders
            files = [f for f in files if not f.startswith('.') and os.path.isfile(os.path.join(INBOX_DIR, f))]
        except Exception as e:
            return {
                "reply": f"❌ Couldn't open inbox: {e}",
                "emotion": "neutral",
                "state": "speaking"
            }

        if not files:
            return {
                "reply": (
                    "📂 Your rexy_inbox/ folder is empty!\n"
                    "Drop any file there and say 'read filename.txt'"
                ),
                "emotion": "neutral",
                "state": "speaking"
            }

        file_list = "\n".join(f"• {f}" for f in files)
        prefix = "Which file should I read? Here's what's in your inbox:\n\n" if hint else "📂 Files in rexy_inbox/:\n\n"
        return {
            "reply": f"{prefix}{file_list}\n\nSay 'read <filename>' to open one.",
            "emotion": "neutral",
            "state": "speaking"
        }

    # ─────────────────────────────────────────────
    # EXTRACT FILENAME FROM MESSAGE
    # ─────────────────────────────────────────────
    def _extract_filename(self, message: str) -> str:
        """
        Try to find a filename in the message.
        "read notes.txt"      → "notes.txt"
        "open my resume.docx" → "resume.docx"
        "read the homework"   → "homework"  (no extension, will fuzzy match)
        """
        # Pattern 1: explicit filename with extension
        match = re.search(r'\b([\w\-]+\.(txt|md|csv|json|pdf|pptx|docx|py|log))\b', message, re.IGNORECASE)
        if match:
            return match.group(1)

        # Pattern 2: "read/open/show X" — grab word after trigger
        match = re.search(
            r'\b(read|open|show|load|display)\s+(?:the\s+)?(?:file\s+)?(?:called\s+)?(?:my\s+)?([\w\-\.]+)',
            message,
            re.IGNORECASE
        )
        if match:
            candidate = match.group(2).strip()
            # Ignore generic words
            if candidate.lower() not in {"file", "files", "inbox", "folder", "it", "this"}:
                return candidate

        return ""

    # ─────────────────────────────────────────────
    # READ FILE
    # ─────────────────────────────────────────────
    def _read_file(self, filename: str) -> Dict[str, Any]:
        """
        Find and read the file from rexy_inbox/.
        Tries exact match first, then fuzzy match (ignoring extension).
        """
        # Exact match
        exact_path = os.path.join(INBOX_DIR, filename)
        if os.path.exists(exact_path):
            return self._parse_and_reply(exact_path, filename)

        # Fuzzy match — find a file whose name starts with or contains filename
        name_no_ext = os.path.splitext(filename)[0].lower()
        try:
            all_files = os.listdir(INBOX_DIR)
        except Exception:
            all_files = []

        for f in all_files:
            if os.path.splitext(f)[0].lower() == name_no_ext:
                path = os.path.join(INBOX_DIR, f)
                return self._parse_and_reply(path, f)

        # Still not found
        return {
            "reply": (
                f"❌ Couldn't find '{filename}' in rexy_inbox/.\n"
                f"Drop the file into rexy_inbox/ first, then try again.\n"
                f"Say 'list inbox' to see what's there."
            ),
            "emotion": "neutral",
            "state": "speaking"
        }

    def _parse_and_reply(self, path: str, filename: str) -> Dict[str, Any]:
        """Read the file and return formatted content based on type."""
        ext = os.path.splitext(filename)[1].lower()

        try:
            if ext in (".txt", ".md", ".log", ".py"):
                content = self._read_text(path)

            elif ext == ".csv":
                content = self._read_csv(path)

            elif ext == ".json":
                content = self._read_json(path)

            elif ext == ".pdf":
                content = self._read_pdf(path)

            elif ext == ".pptx":
                content = self._read_pptx(path)

            elif ext == ".docx":
                content = self._read_docx(path)

            else:
                # Try reading as plain text anyway
                content = self._read_text(path)

        except Exception as e:
            return {
                "reply": f"❌ Error reading '{filename}': {str(e)[:100]}",
                "emotion": "neutral",
                "state": "speaking"
            }

        # Trim if too long
        trimmed = False
        if len(content) > MAX_CHARS:
            content  = content[:MAX_CHARS]
            trimmed  = True

        reply = f"📄 **{filename}**\n\n{content}"
        if trimmed:
            reply += f"\n\n_... (truncated, showing first {MAX_CHARS} characters)_"

        return {
            "reply": reply,
            "emotion": "neutral",
            "state": "speaking"
        }

    # ─────────────────────────────────────────────
    # FORMAT READERS
    # ─────────────────────────────────────────────
    def _read_text(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read().strip()

    def _read_csv(self, path: str) -> str:
        """Show first 20 rows of CSV as plain text table."""
        lines = []
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            for i, line in enumerate(f):
                if i >= 20:
                    lines.append(f"... ({i} more rows)")
                    break
                lines.append(line.rstrip())
        return "\n".join(lines)

    def _read_json(self, path: str) -> str:
        """Pretty print JSON."""
        import json
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)

    def _read_pdf(self, path: str) -> str:
        """Extract text from PDF using pypdf."""
        try:
            from pypdf import PdfReader
        except ImportError:
            return "❌ pypdf not installed. Run: pip install pypdf --break-system-packages"

        reader = PdfReader(path)
        pages  = []
        for i, page in enumerate(reader.pages[:10]):  # Max 10 pages
            text = page.extract_text()
            if text:
                pages.append(f"[Page {i+1}]\n{text.strip()}")
        return "\n\n".join(pages) if pages else "No text found in PDF."

    def _read_pptx(self, path: str) -> str:
        """Extract slide text from PowerPoint."""
        try:
            from pptx import Presentation
        except ImportError:
            return "❌ python-pptx not installed. Run: pip install python-pptx --break-system-packages"

        prs    = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():  # type: ignore
                    texts.append(shape.text.strip())  # type: ignore
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))

        return "\n\n".join(slides) if slides else "No text found in presentation."

    def _read_docx(self, path: str) -> str:
        """Extract paragraphs from Word document."""
        try:
            from docx import Document
        except ImportError:
            return "❌ python-docx not installed. Run: pip install python-docx --break-system-packages"

        doc  = Document(path)
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras) if paras else "No text found in document."

    # ─────────────────────────────────────────────
    # ENSURE INBOX EXISTS
    # ─────────────────────────────────────────────
    def _ensure_inbox(self) -> None:
        """Create rexy_inbox/ if it doesn't exist yet."""
        if not os.path.exists(INBOX_DIR):
            os.makedirs(INBOX_DIR)
            logger.info(f"Created inbox folder: {INBOX_DIR}")